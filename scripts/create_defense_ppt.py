from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

PRIMARY = RGBColor(0x00, 0x5B, 0xAC)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MEDIUM_GRAY = RGBColor(0x99, 0x99, 0x99)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text="", font_size=18,
                font_color=DARK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="微软雅黑"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=18, font_color=DARK,
                  bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(6),
                  font_name="微软雅黑"):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    return p


def set_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text


def add_top_bar(slide):
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), fill_color=PRIMARY)


def add_bottom_bar(slide, page_num, total=16):
    add_shape(slide, Inches(0), Inches(7.1), SLIDE_WIDTH, Inches(0.4), fill_color=PRIMARY)
    add_textbox(slide, Inches(0.5), Inches(7.12), Inches(6), Inches(0.35),
                "中国东北-俄罗斯远东跨境区域SDGs空间支撑条件评估与发展对策",
                font_size=10, font_color=WHITE, font_name="微软雅黑")
    add_textbox(slide, Inches(11), Inches(7.12), Inches(2), Inches(0.35),
                f"{page_num} / {total}", font_size=10, font_color=WHITE,
                alignment=PP_ALIGN.RIGHT, font_name="微软雅黑")


def add_section_title(slide, title_text, subtitle_text=""):
    add_top_bar(slide)
    add_textbox(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.2),
                title_text, font_size=36, font_color=PRIMARY, bold=True,
                alignment=PP_ALIGN.CENTER)
    if subtitle_text:
        add_textbox(slide, Inches(2), Inches(3.8), Inches(9), Inches(0.8),
                    subtitle_text, font_size=20, font_color=GRAY,
                    alignment=PP_ALIGN.CENTER)


def add_content_slide(slide, title, bullets, page_num, notes=""):
    add_top_bar(slide)
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
                title, font_size=28, font_color=PRIMARY, bold=True)
    add_shape(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.03), fill_color=PRIMARY)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(22)
        p.font.color.rgb = DARK
        p.font.name = "微软雅黑"
        p.space_before = Pt(12)
        p.space_after = Pt(6)
        if bullet.startswith("→") or bullet.startswith("•"):
            p.font.color.rgb = GRAY
            p.font.size = Pt(20)

    add_bottom_bar(slide, page_num)
    if notes:
        set_notes(slide, notes)


def add_two_column_slide(slide, title, left_title, left_bullets,
                         right_title, right_bullets, page_num, notes=""):
    add_top_bar(slide)
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
                title, font_size=28, font_color=PRIMARY, bold=True)
    add_shape(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.03), fill_color=PRIMARY)

    add_shape(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(5.3),
              fill_color=LIGHT_BG, line_color=LIGHT_GRAY)
    add_textbox(slide, Inches(1.0), Inches(1.5), Inches(5.1), Inches(0.5),
                left_title, font_size=22, font_color=PRIMARY, bold=True)
    txL = slide.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.1), Inches(4.4))
    tfL = txL.text_frame
    tfL.word_wrap = True
    for i, b in enumerate(left_bullets):
        p = tfL.paragraphs[0] if i == 0 else tfL.add_paragraph()
        p.text = b
        p.font.size = Pt(18)
        p.font.color.rgb = DARK
        p.font.name = "微软雅黑"
        p.space_before = Pt(8)

    add_shape(slide, Inches(6.8), Inches(1.4), Inches(5.5), Inches(5.3),
              fill_color=LIGHT_BG, line_color=LIGHT_GRAY)
    add_textbox(slide, Inches(7.0), Inches(1.5), Inches(5.1), Inches(0.5),
                right_title, font_size=22, font_color=PRIMARY, bold=True)
    txR = slide.shapes.add_textbox(Inches(7.0), Inches(2.1), Inches(5.1), Inches(4.4))
    tfR = txR.text_frame
    tfR.word_wrap = True
    for i, b in enumerate(right_bullets):
        p = tfR.paragraphs[0] if i == 0 else tfR.add_paragraph()
        p.text = b
        p.font.size = Pt(18)
        p.font.color.rgb = DARK
        p.font.name = "微软雅黑"
        p.space_before = Pt(8)

    add_bottom_bar(slide, page_num)
    if notes:
        set_notes(slide, notes)


def add_placeholder_box(slide, left, top, width, height, label):
    shape = add_shape(slide, left, top, width, height,
                      fill_color=RGBColor(0xE8, 0xE8, 0xE8), line_color=MEDIUM_GRAY)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(16)
    p.font.color.rgb = MEDIUM_GRAY
    p.font.name = "微软雅黑"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def create_ppt():
    template_path = "/Users/hbox/Documents/GitHub/毕业/答辩/思源字体【务必安装】/01蓝穹信大，智探风云.pptx"
    prs = Presentation(template_path)
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    total_pages = 16

    # ===== Slide 1: 封面 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15), fill_color=PRIMARY)
    add_shape(slide, Inches(0), Inches(7.0), SLIDE_WIDTH, Inches(0.5), fill_color=PRIMARY)

    add_textbox(slide, Inches(1), Inches(1.0), Inches(11.3), Inches(0.6),
                "南京信息工程大学 地理科学学院", font_size=20, font_color=GRAY,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.5),
                "中国东北-俄罗斯远东跨境区域\nSDGs空间支撑条件评估与发展对策",
                font_size=36, font_color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
    add_shape(slide, Inches(5.5), Inches(3.8), Inches(2.3), Inches(0.04), fill_color=PRIMARY)
    add_textbox(slide, Inches(2), Inches(4.2), Inches(9.3), Inches(0.5),
                "本科毕业论文答辩", font_size=24, font_color=DARK,
                alignment=PP_ALIGN.CENTER)

    info_lines = [
        "答辩人：XXX    学号：XXXXXXXXXX",
        "指导教师：XXX 教授",
        "专    业：自然地理与资源环境 / 人文地理与城乡规划",
        "答辩日期：2026年6月"
    ]
    txBox = slide.shapes.add_textbox(Inches(3), Inches(5.0), Inches(7.3), Inches(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(info_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = GRAY
        p.font.name = "微软雅黑"
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(6)

    set_notes(slide, """【封面页】时长：15秒

开场白：
"各位老师好，我是XX专业XX班的XXX，我的毕业论文题目是《中国东北-俄罗斯远东跨境区域SDGs空间支撑条件评估与发展对策》，指导教师是XXX老师。下面我将从研究背景、数据与方法、主要结果和结论四个方面进行汇报。"

注意事项：
- 站定后先向答辩组老师问好，语速适中
- 论文题目要完整念出，不要省略
- 答辩日期根据实际调整""")

    # ===== Slide 2: 目录 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
                "汇报提纲", font_size=28, font_color=PRIMARY, bold=True)
    add_shape(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.03), fill_color=PRIMARY)

    sections = [
        ("01", "研究背景与意义", "跨境区域SDGs评估的学术缺口与现实需求"),
        ("02", "研究区与数据方法", "四维十二指标体系、10km网格、TOPSIS评价"),
        ("03", "SDG综合指数时空格局", "四期空间分布、国别省州统计"),
        ("04", "ESDA空间集聚特征", "Global Moran's I + LISA集聚类型"),
        ("05", "维度组合与短板识别", "四维得分对比、LISA类型区维度统计"),
        ("06", "跨境差异与协同治理", "中俄互补结构、差异化风险与协同路径"),
        ("07", "结论与展望", "核心发现、不足与展望"),
    ]
    for i, (num, title, desc) in enumerate(sections):
        y = Inches(1.5) + Inches(i * 0.75)
        add_shape(slide, Inches(1.0), y, Inches(0.7), Inches(0.55), fill_color=PRIMARY)
        add_textbox(slide, Inches(1.0), y + Inches(0.05), Inches(0.7), Inches(0.5),
                    num, font_size=22, font_color=WHITE, bold=True,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(2.0), y, Inches(4), Inches(0.35),
                    title, font_size=22, font_color=DARK, bold=True)
        add_textbox(slide, Inches(2.0), y + Inches(0.32), Inches(9), Inches(0.3),
                    desc, font_size=16, font_color=GRAY)

    add_bottom_bar(slide, 2, total_pages)
    set_notes(slide, """【目录页】时长：15秒

过渡语："本次汇报主要包含以下七个部分。"

建议：
- 快速过一遍目录即可，不要展开解释
- 如果答辩组有老师翻论文，目录页能帮助他们对齐进度
- 后续每进入新章节时，可以回顾"现在进入第X部分"起到导航作用""")

    # ===== Slide 3: 研究背景 =====
    add_content_slide(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "01 研究背景",
        [
            "• SDGs本地化评估是联合国2030议程的核心实施路径",
            "• 中俄边境区域发展不平衡、生态与经济矛盾突出",
            "• 跨境区域SDGs评估存在三重缺口：",
            "→ 指标体系缺乏跨境可比性",
            "→ 时空演变分析多停留在单期或两期对比",
            "→ 缺少对边境两侧空间短板差异的系统识别",
            "",
            "• 本研究的定位：",
            "→ 构建跨境可比的四维十二指标SDGs评价体系",
            "→ 实现2010-2024四期时空动态评估",
            "→ 识别跨境空间短板差异，提出分区协同对策"
        ],
        3,
        notes="""【研究背景】时长：约1.5分钟

讲解要点：
1. 先说大背景：SDGs是全球共识，但本地化评估尤其是跨境区域评估不足
2. 再说具体问题：中俄边境区域发展差异大，但缺少系统性的空间评估
3. 最后点明本研究要解决的三个缺口

关键话术：
"SDGs本地化评估是当前可持续发展研究的前沿，但跨境区域的研究几乎空白。中俄边境两侧在生态、经济、基础设施等方面存在显著差异，这种差异缺乏系统性的空间量化评估。"

可能被问：
Q: 你的研究和已有的中俄区域发展研究有什么不同？
A: 已有研究多从经济贸易角度分析，本研究首次从SDGs空间支撑条件角度进行网格尺度的系统评估，覆盖了人类活动、通达性、生态状态、环境本底四个维度。"""
    )

    # ===== Slide 4: 研究意义 =====
    add_two_column_slide(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "01 研究意义",
        "理论意义",
        [
            "• 拓展SDGs评估的跨境空间视角",
            "• 构建适用于高纬度跨境区域的指标体系",
            "• 丰富网格尺度SDGs时空演变分析方法",
            "• 为维度内熵权+维度间等权的组合赋权提供实证"
        ],
        "实践意义",
        [
            "• 为中俄边境协同发展提供空间靶向依据",
            "• 识别跨境互补优势与短板差异",
            "• 支撑边境口岸建设与生态保护协调决策",
            "• 为跨境区域SDGs监测提供可复制的框架"
        ],
        4,
        notes="""【研究意义】时长：约1分钟

讲解要点：
- 理论意义强调"方法创新"：跨境可比指标体系、组合赋权、网格尺度
- 实践意义强调"决策支撑"：空间靶向、互补识别、监测框架
- 不要展开太多，点到为止，详细内容在结果部分展示

过渡语："基于以上研究背景和意义，下面介绍研究区概况与数据方法。"

可能被问：
Q: 你的研究对实际政策有什么影响？
A: 识别了中俄边境两侧的空间短板差异，为跨境互补合作提供了空间靶向依据。例如俄方生态优势与中方基建优势的互补，可以指导边境口岸选址和生态保护红线划定。"""
    )

    # ===== Slide 5: 研究区 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
                "02 研究区概况", font_size=28, font_color=PRIMARY, bold=True)
    add_shape(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.03), fill_color=PRIMARY)

    add_placeholder_box(slide, Inches(0.8), Inches(1.4), Inches(7.5), Inches(5.3),
                        "【插入图2-1 研究区区位图】\n\n含：中俄边界、省州行政区划、\n主要口岸点、地形底图")

    txBox = slide.shapes.add_textbox(Inches(8.8), Inches(1.4), Inches(4), Inches(5.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    items = [
        ("研究区范围", True),
        ("中国东北3省 + 俄罗斯远东5个联邦主体", False),
        ("", False),
        ("核心特征", True),
        ("• 面积约620万km²", False),
        ("• 中方：人口密集、基建完善", False),
        ("• 俄方：地广人稀、生态优良", False),
        ("• 边境线超4300km，口岸20+", False),
        ("", False),
        ("数据口径", True),
        ("• 四期：2010/2015/2020/2024", False),
        ("• 网格：10km（ESRI:102025）", False),
        ("• 行政：GADM 4.1 ADM1", False),
    ]
    for i, (text, is_bold) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(18) if not is_bold else Pt(20)
        p.font.color.rgb = PRIMARY if is_bold else DARK
        p.font.bold = is_bold
        p.font.name = "微软雅黑"
        p.space_before = Pt(4)

    add_bottom_bar(slide, 5, total_pages)
    set_notes(slide, """【研究区概况】时长：约1分钟

讲解要点：
1. 指向地图说明研究区范围：中国东北三省（黑吉辽）+ 俄罗斯远东五个联邦主体
2. 强调核心特征：中俄两侧发展差异显著，这正是跨境评估的价值所在
3. 简要说明数据口径：四期、10km网格、等面积投影

地图替换说明：
- 将QGIS导出的图2-1研究区区位图（高分辨率PNG）替换灰色占位框
- 地图应包含：中俄边界线、省州边界、主要口岸点、地形底图
- 确保地图有图例、比例尺、指北针

过渡语："在这样一个跨境区域，我们构建了以下指标体系和方法。"

可能被问：
Q: 为什么选择这5个俄罗斯联邦主体？
A: 与中国东北接壤的俄罗斯远东联邦主体，包括阿穆尔州、犹太自治州、哈巴罗夫斯克边疆区、滨海边疆区和外贝加尔边疆区，覆盖了中俄边境的主要区域。"""
    )

    # ===== Slide 6: 指标体系 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
                "02 四维十二指标评价体系", font_size=28, font_color=PRIMARY, bold=True)
    add_shape(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.03), fill_color=PRIMARY)

    dims = [
        ("D1 人类活动", PRIMARY, [
            "人口密度 (GPW)", "夜间灯光 (VIIRS)", "耕地比例 (MCD12Q1)"
        ]),
        ("D2 通达性", RGBColor(0x2E, 0x86, 0xC1), [
            "距口岸距离 (负向)", "距城市距离 (负向)", "距道路距离 (负向)"
        ]),
        ("D3 生态状态", RGBColor(0x1E, 0x8A, 0x49), [
            "NDVI (MOD13A3)", "林地覆盖率 (MCD12Q1)", "净初级生产力 (MOD17)"
        ]),
        ("D4 环境本底", RGBColor(0xD4, 0x7A, 0x0A), [
            "地形坡度 (负向, GTOPO30)", "气候水分亏缺 (负向, TerraClimate)", "低温冷胁迫 (负向, TerraClimate)"
        ]),
    ]

    for i, (dim_name, color, indicators) in enumerate(dims):
        x = Inches(0.8) + Inches(i * 3.1)
        add_shape(slide, x, Inches(1.5), Inches(2.9), Inches(0.6), fill_color=color)
        add_textbox(slide, x + Inches(0.1), Inches(1.52), Inches(2.7), Inches(0.55),
                    dim_name, font_size=18, font_color=WHITE, bold=True,
                    alignment=PP_ALIGN.CENTER)
        for j, ind in enumerate(indicators):
            y = Inches(2.3) + Inches(j * 1.0)
            add_shape(slide, x, y, Inches(2.9), Inches(0.85),
                      fill_color=LIGHT_BG, line_color=LIGHT_GRAY)
            add_textbox(slide, x + Inches(0.15), y + Inches(0.08), Inches(2.6), Inches(0.7),
                        ind, font_size=16, font_color=DARK, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.5),
                "赋权方法：维度内熵权法（按离散信息量分配） + 维度间等权（避免单一维度主导）",
                font_size=18, font_color=GRAY, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.5),
                "综合评价：TOPSIS相对贴近度 → SDG综合指数",
                font_size=20, font_color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

    add_bottom_bar(slide, 6, total_pages)
    set_notes(slide, """【指标体系】时长：约1.5分钟

讲解要点：
1. 四个维度覆盖SDGs的核心面向：人类活动、通达性、生态状态、环境本底
2. 每个维度三个指标，正向指标越大越好，负向指标反向处理
3. 强调赋权逻辑：维度内用熵权（尊重数据离散信息），维度间用等权（保证公平）
4. TOPSIS计算相对贴近度作为综合指数

关键话术：
"指标体系涵盖四个维度十二个指标。维度内部采用熵权法，让数据自身说话；维度之间采用等权，避免高离散维度过度主导综合指数。最终用TOPSIS相对贴近度作为SDG综合指数。"

年份口径说明（如果被问）：
- 2010模型年夜光采用2012年VIIRS近似替代
- 2024年人口沿用2020基线
- 坡度为静态DEM指标，所有年份共用

可能被问：
Q: 为什么维度间用等权而不是全熵权？
A: 四个维度代表SDGs的不同面向，全熵权会让高离散维度（如D1人类活动）过度主导综合指数，等权保证各维度公平贡献。我们也做了加权求和对比验证，结果高度相关。"""
    )

    # ===== Slide 7: 技术路线 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
                "02 技术路线", font_size=28, font_color=PRIMARY, bold=True)
    add_shape(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.03), fill_color=PRIMARY)

    add_placeholder_box(slide, Inches(2.5), Inches(1.4), Inches(8.3), Inches(5.3),
                        "【插入图3-1 技术路线图】\n\n含：数据获取→指标计算→标准化→赋权\n→TOPSIS→趋势分析→热点识别→LISA→分区对策")

    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(1.5), Inches(5.0),
                "数据获取\n↓\n指标计算\n↓\n极差标准化\n↓\n熵权+等权\n↓\nTOPSIS\n↓\n趋势分析\n↓\nLISA+短板\n↓\n分区对策",
                font_size=16, font_color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

    add_bottom_bar(slide, 7, total_pages)
    set_notes(slide, """【技术路线】时长：约1分钟

讲解要点：
1. 按流程图从上到下讲解，强调方法链的完整性
2. 关键步骤：全时期合并极差标准化 → 逐年维度内熵权 → 维度间等权 → TOPSIS
3. 时空分析：变化率统计 + ESDA空间集聚 + 维度组合与短板识别
4. 最终产出：跨境差异解释与协同治理路径

替换说明：
- 将论文中的图3-1技术路线图（PPTX可编辑版）替换灰色占位框
- 如果没有现成技术路线图，可以用PPT自带的SmartArt或流程图工具制作

过渡语："基于以上方法，下面展示核心研究结果。"

可能被问：
Q: 为什么用TOPSIS而不是简单加权求和？
A: TOPSIS通过正负理想解计算相对贴近度，能更好反映各方案与最优/最劣方案的相对位置，比简单加权更稳健。我们也做了加权求和对比验证，相关系数>0.99，说明结果稳健。

Q: 全时期合并标准化有什么好处？
A: 保证不同年份的指标在同一尺度上可比，如果逐年标准化则无法进行跨年比较。"""
    )

    # ===== Slide 8: 四期SDG指数空间格局 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
                "03 SDG综合指数四期空间格局", font_size=28, font_color=PRIMARY, bold=True)
    add_shape(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.03), fill_color=PRIMARY)

    for i, year in enumerate([2010, 2015, 2020, 2024]):
        x = Inches(0.5) + Inches(i * 3.15)
        add_placeholder_box(slide, x, Inches(1.4), Inches(2.95), Inches(4.5),
                            f"【插入{year}年\nSDG综合指数空间分布图】")
        add_textbox(slide, x, Inches(6.0), Inches(2.95), Inches(0.4),
                    str(year), font_size=20, font_color=PRIMARY, bold=True,
                    alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4),
                '核心发现：SDG综合指数呈"东高西低、南高北低"格局，中方整体高于俄方，2024年较2010年整体提升',
                font_size=18, font_color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    add_bottom_bar(slide, 8, total_pages)
    set_notes(slide, """【四期空间格局】时长：约2.5分钟 ★核心页面

讲解要点：
1. 逐期展示四幅地图，说明空间格局特征
2. 强调"东高西低、南高北低"的总体格局
3. 指出中方整体高于俄方的原因：人口密度、基建完善度、经济活力
4. 2024年较2010年整体提升，但提升幅度空间不均

关键话术：
"从四期SDG综合指数空间分布来看，整体呈'东高西低、南高北低'格局。中方区域受人口集聚和基建完善驱动，指数普遍较高；俄方远东地区地广人稀，D1人类活动维度得分偏低，但D3生态状态维度表现优异。2024年较2010年整体有所提升，但提升幅度存在明显空间差异。"

地图替换说明：
- 将QGIS导出的四期SDG综合指数空间分布图替换灰色占位框
- 四幅地图应使用统一的色阶和分类断点
- 确保每幅地图有图例、比例尺

可能被问：
Q: 为什么俄方指数这么低？
A: 俄方远东地区人口密度极低（部分网格<1人/km²），基建覆盖率低，导致D1人类活动和D2通达性维度得分偏低。但D3生态状态维度表现优异，体现了"生态优势-发展短板"的跨境互补特征。"""
    )

    # ===== Slide 9: 国别省州统计 =====
    add_two_column_slide(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "03 国别与省州级统计",
        "中俄对比",
        [
            "• 中方均值：0.XXX（2024年）",
            "• 俄方均值：0.XXX（2024年）",
            "• 中方优势维度：D1人类活动、D2通达性",
            "• 俄方优势维度：D3生态状态",
            "• D4环境本底双方均较严峻",
            "",
            "【替换为实际统计数据表/柱状图】"
        ],
        "省州排名（2024年）",
        [
            "• 中方前三：XXX、XXX、XXX",
            "• 中方后三：XXX、XXX、XXX",
            "• 俄方前三：XXX、XXX、XXX",
            "• 俄方后三：XXX、XXX、XXX",
            "",
            "【替换为省州维度得分柱状图】",
            "【数据来源：admin1_group_statistics.csv】"
        ],
        9,
        notes="""【国别省州统计】时长：约1.5分钟

讲解要点：
1. 先说中俄整体对比，突出互补特征
2. 再说省州排名，指出最高和最低的区域
3. 用具体数字说话，不要只说"高""低"

数据替换说明：
- 将admin1_group_statistics.csv中的实际数据填入
- 建议用分组柱状图展示四期各维度得分
- 数据来源：data/admin1_group_statistics.csv, data/country_group_statistics.csv

关键话术：
"从国别统计来看，中方在D1人类活动和D2通达性维度具有明显优势，俄方在D3生态状态维度表现突出，体现了典型的跨境互补特征。"

可能被问：
Q: 省州差异的主要原因是什么？
A: 省会城市和口岸城市周边的网格在D1人类活动和D2通达性维度得分显著高于偏远地区，这与人口分布和基建覆盖高度相关。"""
    )

    # ===== Slide 10: 变化趋势 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
                "04 ESDA空间集聚特征", font_size=28, font_color=PRIMARY, bold=True)
    add_shape(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.03), fill_color=PRIMARY)

    add_placeholder_box(slide, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.0),
                        "【插入2024年SDG综合指数\nLISA集聚图】\n\n高-高/低-低/高-低/低-高")

    add_placeholder_box(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.0),
                        "【插入2010-2024变化率\nLISA集聚图】")

    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(6.0), Inches(0.4),
                "2024年综合指数LISA", font_size=16, font_color=PRIMARY, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.8), Inches(6.5), Inches(6.0), Inches(0.4),
                "变化率LISA", font_size=16, font_color=PRIMARY, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_bottom_bar(slide, 10, total_pages)
    set_notes(slide, """【ESDA空间集聚特征】时长：约2分钟

讲解要点：
1. 左图：2024年综合指数LISA集聚图——高-高集聚区在哪、低-低集聚区在哪
2. 右图：变化率LISA集聚图——哪些区域的变化存在空间集聚
3. 先报全局Moran's I值，说明存在显著空间正相关

关键话术：
"全局Moran's I为0.9179（2024年），通过1%显著性检验，表明SDG综合指数存在显著空间正相关。LISA集聚图显示，高-高集聚区主要分布在中方省会城市周边，低-低集聚区主要分布在俄方远东内陆。变化率LISA显示，提升显著区域同样存在空间集聚。"

地图替换说明：
- 左图：esda_lisa_sdg_index_2024.png
- 右图：esda_lisa_change_rate_2010_2024.png
- 全局Moran's I值：data/esda_global_moran.csv
- LISA统计：data/esda_lisa_statistics.csv

可能被问：
Q: 为什么用LISA而不是Gi*？
A: LISA能同时识别四种集聚类型（高-高、低-低、高-低、低-高），比Gi*提供更丰富的空间异质性信息，更适合识别跨境区域的集聚-孤立模式。""")

    # ===== Slide 11: LISA集聚 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
                "05 维度组合与短板识别", font_size=28, font_color=PRIMARY, bold=True)
    add_shape(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.03), fill_color=PRIMARY)

    add_placeholder_box(slide, Inches(0.5), Inches(1.4), Inches(6.0), Inches(4.5),
                        "【插入四维得分对比图】\n\n中俄两侧D1-D4维度\n均值对比柱状图/雷达图")

    add_placeholder_box(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(4.5),
                        "【插入LISA类型区维度热力图】\n\nHH/LL/HL/LH各维度\n均值热力图")

    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(6.0), Inches(0.4),
                "中俄维度得分对比", font_size=16, font_color=PRIMARY, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.8), Inches(6.0), Inches(6.0), Inches(0.4),
                "LISA类型区维度热力图", font_size=16, font_color=PRIMARY, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4),
                "两侧最强维度均为D2通达性，短板维度均为D1人类活动；HH与LL差值最大维度为D3生态状态",
                font_size=18, font_color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    add_bottom_bar(slide, 11, total_pages)
    set_notes(slide, """【维度组合与短板识别】时长：约1.5分钟

讲解要点：
1. 左图：中俄两侧D1-D4维度得分对比——两侧最强维度均为D2通达性，短板维度均为D1人类活动
2. 右图：LISA类型区维度热力图——HH与LL集聚区差值最大维度为D3生态状态（0.5906），其次为D2通达性（0.3650）
3. 核心发现：中方D1人类活动和D2通达性显著高于俄方，俄方D3生态状态优于中方

数据替换说明：
- 左图：data/country_group_statistics.csv
- 右图：data/lisa_dimension_shortboard_statistics.csv
- 建议用分组柱状图或雷达图展示维度对比

关键话术：
"维度组合分析显示，两侧最强维度均为D2通达性，短板维度均为D1人类活动。HH与LL集聚区差值最大维度为D3生态状态，体现了中俄跨境的生态-发展互补结构。"

可能被问：
Q: 短板识别的可靠性如何？
A: 基于LISA显著网格（p<0.05）进行统计，样本量充足。同时我们做了指标层面的短板分析，结果与维度层面一致。""")

    # ===== Slide 12: 短板识别 =====
    add_content_slide(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "05 指标层面短板详情",
        [
            "• HH集聚区（高-高）：主短板指标【替换为实际数据】",
            "• LL集聚区（低-低）：主短板指标【替换为实际数据】",
            "• HL集聚区（高-低）：主短板指标【替换为实际数据】",
            "• LH集聚区（低-高）：主短板指标【替换为实际数据】",
            "",
            "• 核心发现：",
            "→ 中方低值区主要受D4环境本底制约（坡度、冷胁迫）",
            "→ 俄方低值区主要受D1人类活动和D2通达性制约",
            "",
            "【替换为指标短板统计表/热力图】"
        ],
        12,
        notes="""【指标层面短板详情】时长：约1分钟

讲解要点：
1. 按LISA集聚类型逐一说明主短板指标
2. 核心发现：中方低值区受环境本底制约，俄方低值区受人类活动制约

数据替换说明：
- 短板统计：data/lisa_dimension_shortboard_statistics.csv
- 指标短板详情：data/lisa_indicator_shortboard_details.csv

过渡语："基于以上维度组合和短板分析，下面讨论跨境差异与协同治理路径。"

可能被问：
Q: 短板识别对政策有什么启示？
A: 短板识别揭示了不同类型区域的差异化制约因素，为精准施策提供了空间靶向依据。"""
    )

    # ===== Slide 13: 跨境差异与分区对策 =====
    add_two_column_slide(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "06 跨境差异与协同治理路径",
        "跨境核心差异",
        [
            "• 中方优势：人口集聚、基建完善、口岸联通",
            "• 俄方优势：生态优良、资源丰富、林地覆盖高",
            "• 共同挑战：气候水分亏缺、高纬度冷胁迫",
            "",
            "• 互补结构：",
            '→ 中方"发展优势+生态短板"',
            '→ 俄方"生态优势+发展短板"',
            "→ 口岸节点是互补关键枢纽"
        ],
        "协同治理路径",
        [
            "① 强化口岸节点枢纽功能",
            "  → 提升D2通达性，促进要素流动",
            "② 生态-发展互补合作",
            "  → 中方基建经验+俄方生态资源",
            "③ 差异化风险应对",
            "  → 中方关注生态约束，俄方关注人口流失",
            "④ 共同应对气候胁迫",
            "  → D4环境本底是双方共同挑战"
        ],
        13,
        notes="""【跨境差异与协同治理】时长：约1.5分钟

讲解要点：
1. 左栏：概括中俄核心差异和互补特征
2. 右栏：四条协同治理路径，每条对应一个维度
3. 强调路径与空间格局和短板识别的对应关系

关键话术：
"基于维度组合和短板分析，我们提出四条协同治理路径：强化口岸枢纽、生态-发展互补、差异化风险应对、共同应对气候胁迫。这些路径与空间格局和短板识别结果直接对应。"

可能被问：
Q: 你的治理路径和已有研究有什么不同？
A: 已有对策多基于定性分析，本研究基于网格尺度的定量评估和LISA集聚类型识别，路径具有空间靶向性。

Q: 这些路径的可操作性如何？
A: 路径与具体的省州和集聚类型对应，可以指导口岸选址优化、生态保护红线划定、基建投资优先区域等实际决策。"""
    )

    # ===== Slide 14: 结论 =====
    add_content_slide(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "07 主要结论",
        [
            "① 构建了跨境可比的四维十二指标SDGs评价体系，",
            "   实现了中俄东北-远东区域2010-2024四期时空评估",
            "",
            '② SDG综合指数呈"东高西低、南高北低"格局，',
            "   中方整体高于俄方，2024年较2010年整体提升",
            "",
            "③ ESDA分析揭示显著空间集聚（Moran's I=0.9179），",
            "   HH集聚区集中在中方省会周边，LL集聚区在俄方内陆",
            "",
            "④ 维度组合分析揭示跨境互补结构：",
            "   两侧最强维度为D2通达性，短板维度为D1人类活动",
            "",
            "⑤ 提出四条协同治理路径，为跨境SDGs协调提供空间靶向依据"
        ],
        14,
        notes="""【主要结论】时长：约1分钟

讲解要点：
1. 五条结论逐一念出，每条一句话
2. 结论要呼应研究问题，形成闭环
3. 不要展开解释，结论页要简洁有力

关键话术：
"综上所述，本研究构建了跨境可比的SDGs评价体系，揭示了中俄边境区域的空间格局、变化趋势和短板差异，提出了差异化的协同分区对策。"

过渡语："下面简要说明本研究的不足与展望。"

可能被问：
Q: 你觉得最重要的发现是什么？
A: 最重要的发现是跨境互补结构——中方发展优势与俄方生态优势的互补，这为边境协同发展提供了明确的空间靶向。"""
    )

    # ===== Slide 15: 不足与展望 =====
    add_content_slide(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "不足与展望",
        [
            "研究不足：",
            "• 2024年人口指标沿用2020基线，未反映最新人口变化",
            "• 2010模型年夜光采用2012年VIIRS近似替代",
            "• 10km网格尺度可能掩盖更细粒度的空间异质性",
            "• 未纳入政策制度等定性维度的量化指标",
            "",
            "未来展望：",
            "• 引入WorldPop等更高分辨率人口数据",
            "• 结合DMSP-OLS与VIIRS构建跨传感器夜光序列",
            "• 探索5km或1km网格尺度的敏感性分析",
            "• 纳入政策文本量化指标，完善评价体系"
        ],
        15,
        notes="""【不足与展望】时长：约30秒

讲解要点：
1. 诚实说明不足，但不要过度自我否定
2. 每条不足配一个改进方向
3. 展望要具体、可行，不要空泛

关键话术：
"本研究存在以下不足：2024年人口沿用2020基线、2010年夜光采用2012年近似替代、10km网格可能掩盖更细粒度差异。未来可引入更高分辨率数据、构建跨传感器夜光序列、探索更细网格尺度。"

⚠️ 注意：
- 不足要主动说，不要等老师追问
- 但语气要平和，不要说"本研究很差"
- 展望要与不足一一对应"""
    )

    # ===== Slide 16: 致谢 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15), fill_color=PRIMARY)
    add_shape(slide, Inches(0), Inches(7.0), SLIDE_WIDTH, Inches(0.5), fill_color=PRIMARY)

    add_textbox(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(1.0),
                "恳请各位老师批评指正", font_size=36, font_color=PRIMARY, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_shape(slide, Inches(5.5), Inches(3.7), Inches(2.3), Inches(0.04), fill_color=PRIMARY)

    add_textbox(slide, Inches(2), Inches(4.2), Inches(9.3), Inches(0.5),
                "答辩人：XXX    指导教师：XXX 教授", font_size=22, font_color=GRAY,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2), Inches(4.8), Inches(9.3), Inches(0.5),
                "南京信息工程大学 地理科学学院", font_size=20, font_color=GRAY,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2), Inches(5.3), Inches(9.3), Inches(0.5),
                "2026年6月", font_size=18, font_color=MEDIUM_GRAY,
                alignment=PP_ALIGN.CENTER)

    set_notes(slide, """【致谢页】时长：约15秒

结尾话术：
"以上是我的汇报内容。本研究的主要不足在于：2024年人口指标沿用2020基线、2010年夜光采用2012年近似替代，以及10km网格尺度可能掩盖更细粒度的空间异质性。恳请各位老师批评指正。"

⚠️ 注意：
- 不要说"感谢聆听"，用"恳请各位老师批评指正"
- 致谢页配色与封面保持一致，首尾呼应
- 停顿等待老师提问

========================================
答辩高频问题预判与应答（备忘）
========================================

Q1: 为什么用TOPSIS而不是简单加权求和？
A: TOPSIS通过正负理想解计算相对贴近度，能更好反映各方案与最优/最劣方案的相对位置，比简单加权更稳健。我们也做了加权求和对比验证，结果高度相关。

Q2: 2010年夜光为什么用2012年数据？
A: VIIRS DNB卫星2012年4月才开始运行，2010年无实测VIIRS数据。采用2012年作为起始期近似替代，论文中已明确披露。

Q3: 为什么用LISA而不是Gi*？
A: LISA能同时识别四种集聚类型（高-高、低-低、高-低、低-高），比Gi*提供更丰富的空间异质性信息，更适合识别跨境区域的集聚-孤立模式。

Q4: 10km网格尺度的选择依据？
A: 10km是跨境区域尺度研究的常用网格大小，既能捕捉省州级空间差异，又不会因网格过小导致数据稀疏。

Q5: 维度间为什么用等权而不是全熵权？
A: 四个维度代表SDGs的不同面向，全熵权会让高离散维度（如D1人类活动）过度主导综合指数，等权保证各维度公平贡献。

Q6: 你的研究有什么实际政策意义？
A: 识别了中俄边境两侧的空间短板差异，为跨境互补合作提供了空间靶向依据，如俄方生态优势与中方基建优势的互补。

Q7: 全时期合并标准化有什么好处？
A: 保证不同年份的指标在同一尺度上可比，如果逐年标准化则无法进行跨年比较。

Q8: 为什么选择这5个俄罗斯联邦主体？
A: 与中国东北接壤的俄罗斯远东联邦主体，覆盖了中俄边境的主要区域。

Q9: 你的研究和已有的中俄区域发展研究有什么不同？
A: 已有研究多从经济贸易角度分析，本研究首次从SDGs空间支撑条件角度进行网格尺度的系统评估。

Q10: 数据质量如何保证？
A: 所有数据来自权威来源（GPW、VIIRS、MODIS、TerraClimate、GTOPO30），年份映射和近似替代已在论文中明确披露。

========================================
评优硬性指标备忘
========================================
- 综合成绩≥90分（指导40%+评阅30%+答辩30%）
- 查重率＜10%（维普，答辩门槛＜20%）
- AIGC检测＜40%
- 校级评优名额：毕业生3%
- 一二等奖需成果支撑：论文/竞赛/专利/软著
- 汇报时间：10-15分钟
- 提问时间：5-10分钟""")

    output_path = "/Users/hbox/Documents/GitHub/毕业/答辩PPT模板_南信大地理科学学院.pptx"
    prs.save(output_path)
    print(f"PPT saved to: {output_path}")


if __name__ == "__main__":
    create_ppt()
